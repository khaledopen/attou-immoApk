import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_soutenance_presentation():
    prs = Presentation()
    
    # Dimensions (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Couleurs du thème AttouHome
    c_dark = RGBColor(15, 23, 42)      # #0f172a (Bleu nuit profond)
    c_sky = RGBColor(14, 165, 233)     # #0ea5e9 (Bleu ciel)
    c_orange = RGBColor(249, 115, 22)  # #f97316 (Orange vibrant)
    c_red = RGBColor(239, 68, 68)      # #ef4444 (Rouge Alerte)
    c_gray_dark = RGBColor(51, 65, 85) # #334155 (Slate foncé)
    c_gray = RGBColor(100, 116, 139)   # #64748b (Gris moyen)
    c_light = RGBColor(248, 250, 252)  # #f8fafc (Gris très clair)
    c_white = RGBColor(255, 255, 255)  # Blanc
    c_card_bg = RGBColor(255, 255, 255)# Blanc pour les cartes
    
    blank_layout = prs.slide_layouts[6]
    
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, subtitle_text=None, is_dark=False):
        # Bandeau décoratif en haut
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.1))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = c_orange
        stripe.line.color.rgb = c_orange
        
        # Boîte de titre
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.3))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.margin_left = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c_white if is_dark else c_dark
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Segoe UI'
            p2.font.size = Pt(13)
            p2.font.italic = True
            p2.font.color.rgb = c_sky
            p2.space_before = Pt(4)

    def add_footer(slide, current_slide, total_slides, is_dark=False):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
        tf = footer_box.text_frame
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.margin_left = Inches(0)
        p = tf.paragraphs[0]
        p.text = f"AttouHome — Soutenance de Projet  |  {current_slide} / {total_slides}"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(10)
        p.font.color.rgb = c_gray if not is_dark else c_sky
        p.alignment = PP_ALIGN.RIGHT

    def create_card(slide, left, top, width, height, bg_color, border_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.color.rgb = bg_color
        return card

    total_slides = 11

    # ==================== SLIDE 1 : PAGE DE GARDE (DARK) ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, c_dark)
    
    # Accent color block on left
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_orange
    accent_bar.line.color.rgb = c_orange

    # Title & Subtitle box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_brand = tf.paragraphs[0]
    p_brand.text = "PROJET ATTOUHOME"
    p_brand.font.name = 'Segoe UI'
    p_brand.font.size = Pt(16)
    p_brand.font.bold = True
    p_brand.font.color.rgb = c_orange
    p_brand.space_after = Pt(12)
    
    p_title = tf.add_paragraph()
    p_title.text = "Soutenance de Fin de Projet"
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    
    p_subtitle = tf.add_paragraph()
    p_subtitle.text = "Plateforme Innovante de Location Immobilière Assistée en Côte d'Ivoire"
    p_subtitle.font.name = 'Segoe UI'
    p_subtitle.font.size = Pt(20)
    p_subtitle.font.color.rgb = c_sky
    p_subtitle.space_before = Pt(10)
    p_subtitle.space_after = Pt(40)
    
    p_meta = tf.add_paragraph()
    p_meta.text = "Présenté aujourd'hui aux Responsables  •  Démo & Bilan Technique"
    p_meta.font.name = 'Segoe UI'
    p_meta.font.size = Pt(13)
    p_meta.font.italic = True
    p_meta.font.color.rgb = c_white
    p_meta.space_before = Pt(20)

    add_footer(slide1, 1, total_slides, is_dark=True)

    # ==================== SLIDE 2 : CONTEXTE & PROBLÉMATIQUE (LIGHT) ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, c_light)
    add_header(slide2, "Contexte & Problématique", "Les défis de la location immobilière en Côte d'Ivoire (Abidjan)")
    
    # Left Card: Les Locataires
    create_card(slide2, 0.8, 1.8, 5.6, 4.8, c_card_bg, c_sky)
    left_box = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p_l = ltf.paragraphs[0]
    p_l.text = "DÉFIS CÔTÉ LOCATAIRES (TENANTS)"
    p_l.font.name = 'Segoe UI'
    p_l.font.size = Pt(16)
    p_l.font.bold = True
    p_l.font.color.rgb = c_dark
    p_l.space_after = Pt(14)
    
    bullets_l = [
        "Recherche complexe et non ciblée géographiquement (adresses floues).",
        "Difficulté à estimer la distance et le temps de trajet vers les biens.",
        "Manque de traçabilité et de sérieux dans les échanges avec les bailleurs.",
        "Insécurité face aux annonces frauduleuses ou suspectes.",
        "Visites chronophages organisées sans cadre de rendez-vous strict."
    ]
    for b in bullets_l:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    # Right Card: Les Propriétaires
    create_card(slide2, 6.9, 1.8, 5.6, 4.8, c_card_bg, c_orange)
    right_box = slide2.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "DÉFIS CÔTÉ PROPRIÉTAIRES (OWNERS)"
    p_r.font.name = 'Segoe UI'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_dark
    p_r.space_after = Pt(14)
    
    bullets_r = [
        "Saisie d'adresses imprécises sans outil d'autocomplétion cartographique.",
        "Surcharge de sollicitations informelles par messagerie directe.",
        "Absence de planification structurée pour les visites des biens.",
        "Difficulté à gérer la conformité et la modération des annonces.",
        "Risque de saturation avec des demandes obsolètes ou restées sans réponse."
    ]
    for b in bullets_r:
        p_b = rtf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    add_footer(slide2, 2, total_slides)

    # ==================== SLIDE 3 : LA SOLUTION ATTOUHOME (LIGHT) ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, c_light)
    add_header(slide3, "La Solution AttouHome", "Une mise en relation directe, fluide et sécurisée")
    
    # 3 Columns Layout
    col_width = 3.6
    gap = 0.4
    start_left = 0.8
    
    solutions = [
        {
            "title": "1. GÉOLOCALISATION APPLIQUÉE",
            "border": c_sky,
            "items": [
                "Recherche sur carte interactive avec positionnement GPS.",
                "Calcul d'itinéraire routier en temps réel via l'API OSRM.",
                "Sélecteur d'adresse intelligent (Nominatim) limité à la Côte d'Ivoire."
            ]
        },
        {
            "title": "2. SÉCURISATION DU CHAT",
            "border": c_orange,
            "items": [
                "Ouverture de discussion conditionnée à une visite acceptée.",
                "Bandeau d'alerte et de statut pour les deux parties.",
                "Notifications sonores immersives lors de la réception de messages."
            ]
        },
        {
            "title": "3. PROCESSUS AUTOMATISÉS",
            "border": c_dark,
            "items": [
                "Expiration automatique des visites sous 72h avec alerte temps réel.",
                "Modération stricte : signalement d'annonces suspectes en 24h.",
                "Désactivation automatique du chat en cas de modération active."
            ]
        }
    ]
    
    for i, sol in enumerate(solutions):
        current_left = start_left + i * (col_width + gap)
        create_card(slide3, current_left, 1.8, col_width, 4.8, c_card_bg, sol["border"])
        box = slide3.shapes.add_textbox(Inches(current_left + 0.2), Inches(2.0), Inches(col_width - 0.4), Inches(4.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = sol["title"]
        p.font.name = 'Segoe UI'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = sol["border"]
        p.space_after = Pt(15)
        
        for item in sol["items"]:
            p_b = tf.add_paragraph()
            p_b.text = "✔ " + item
            p_b.font.name = 'Segoe UI'
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = c_gray_dark
            p_b.space_before = Pt(10)
            
    add_footer(slide3, 3, total_slides)

    # ==================== SLIDE 4 : ARCHITECTURE TECHNIQUE (LIGHT) ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, c_light)
    add_header(slide4, "Architecture Technique & Stack", "Une architecture à trois tiers moderne, réactive et typée")
    
    # Left Card: Frontend (Client Layers)
    create_card(slide4, 0.8, 1.8, 5.6, 4.8, c_card_bg, c_dark)
    left_box = slide4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p_l = ltf.paragraphs[0]
    p_l.text = "FRONTEND (MULTI-PLATEFORME)"
    p_l.font.name = 'Segoe UI'
    p_l.font.size = Pt(16)
    p_l.font.bold = True
    p_l.font.color.rgb = c_dark
    p_l.space_after = Pt(12)
    
    front_points = [
        "React Native (Expo SDK 54, Expo Router, TypeScript) pour les applications mobiles.",
        "Tenant App : Recherche cartographique, demandes de visite et messagerie.",
        "Owner App : Publication d'annonces, planification et validation des visites.",
        "Admin Web : Interface de modération en React pour les administrateurs.",
        "Modules natifs gérés : react-native-maps, expo-av (audio), expo-location (GPS)."
    ]
    for pt in front_points:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + pt
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(12.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    # Right Card: Backend & Database
    create_card(slide4, 6.9, 1.8, 5.6, 4.8, c_card_bg, c_sky)
    right_box = slide4.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "BACKEND & BASE DE DONNÉES"
    p_r.font.name = 'Segoe UI'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_sky
    p_r.space_after = Pt(12)
    
    back_points = [
        "Node.js + Express API en TypeScript, assurant un typage de bout en bout.",
        "Socket.io pour les communications bidirectionnelles et instantanées.",
        "Prisma ORM pour la modélisation de données et les requêtes sécurisées.",
        "PostgreSQL pour la persistance relationnelle.",
        "CRON / Interval Workers pour les tâches asynchrones récurrentes.",
        "APIs tiers légères : OSRM (itinéraires) et Nominatim (géocodage 🇨🇮)."
    ]
    for pt in back_points:
        p_b = rtf.add_paragraph()
        p_b.text = "• " + pt
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(12.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    add_footer(slide4, 4, total_slides)

    # ==================== SLIDE 5 : MODÈLE DE DONNÉES (LIGHT) ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, c_light)
    add_header(slide5, "Modèle de Données & Relations", "Un schéma relationnel robuste conçu avec Prisma ORM")
    
    # 4 grid cards
    card_w = 5.6
    card_h = 2.2
    
    # Card 1: User
    create_card(slide5, 0.8, 1.8, card_w, card_h, c_card_bg, c_sky)
    box = slide5.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(card_w - 0.2), Inches(card_h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Utilisateur (User)"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = c_sky
    p.space_after = Pt(4)
    bullets = ["Rôles distincts : PROPRIETAIRE et LOCATAIRE", "Hachage Bcrypt des mots de passe", "Authentification basée sur des tokens JWT"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Card 2: Property/Annonce
    create_card(slide5, 6.9, 1.8, card_w, card_h, c_card_bg, c_orange)
    box = slide5.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(card_w - 0.2), Inches(card_h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bien & Annonce (Property / Annonce)"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = c_orange
    p.space_after = Pt(4)
    bullets = ["Description complète, prix, surface, nombre de chambres", "Coordonnées GPS et adresse physique (Côte d'Ivoire)", "Statuts de publication : PUBLIEE, SUSPENDUE, BROUILLON"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Card 3: Visite (Visit)
    create_card(slide5, 0.8, 4.4, card_w, card_h, c_card_bg, c_dark)
    box = slide5.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(card_w - 0.2), Inches(card_h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Demande de Visite (Visit)"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = c_dark
    p.space_after = Pt(4)
    bullets = ["Stocke la date, l'heure et l'identifiant du bien", "États du flux : EN_ATTENTE, ACCEPTEE, REFUSEE, EXPIREE", "Relation un-à-plusieurs (Un bien peut recevoir plusieurs demandes)"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Card 4: Discussion (Chat)
    create_card(slide5, 6.9, 4.4, card_w, card_h, c_card_bg, c_red)
    box = slide5.shapes.add_textbox(Inches(7.0), Inches(4.5), Inches(card_w - 0.2), Inches(card_h - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Discussion & Messages (Chat)"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(4)
    bullets = ["Salon de chat directement associé à un bien immobilier", "Condition : Visite 'ACCEPTEE' obligatoire pour lancer le chat", "Persistance en BDD des messages avec horodatage strict"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(11)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    add_footer(slide5, 5, total_slides)

    # ==================== SLIDE 6 : CARTOGRAPHIE & ITINÉRAIRE (LIGHT) ====================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, c_light)
    add_header(slide6, "Cartographie & Itinéraire en Temps Réel", "Une expérience de navigation fluide pour cibler précisément les logements")
    
    left_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "FONCTIONNALITÉS TECHNIQUES"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_sky
    p.space_after = Pt(14)
    bullets = [
        "Permissions GPS d'expo-location demandées à l'ouverture de la carte.",
        "Autocomplétion des adresses via Nominatim limitée à la Côte d'Ivoire.",
        "Positionnement intelligent : si la latitude/longitude du bien est vide, l'algorithme utilise la commune avec un 'Jitter' (décalage minime) pour éviter que les épingles se superposent sur la carte.",
        "Redirection directe 'Voir sur la carte' depuis la fiche du bien avec passage d'argument focus (id) et zoom immédiat."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    right_box = slide6.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "DOUBLON DE CALCUL D'ITINÉRAIRE"
    p_r.font.name = 'Segoe UI'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_orange
    p_r.space_after = Pt(14)
    solutions = [
        "Requête automatique à OSRM (Open Source Routing Machine) pour tracer la route réelle avec distance/temps de trajet (ex: '🚗 4.2 km (8 min)').",
        "Algorithme de Fallback : si l'API OSRM est indisponible ou l'appareil est hors-ligne, calcul mathématique géométrique de la distance via la Formule de Haversine.",
        "Affichage instantané d'une ligne directe avec l'indicateur '(direct)' pour garantir une expérience sans crash (ex: '5.1 km (direct)')."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Segoe UI'
        p_s.font.size = Pt(13)
        p_s.font.color.rgb = c_gray_dark
        p_s.space_before = Pt(8)

    add_footer(slide6, 6, total_slides)

    # ==================== SLIDE 7 : MESSAGERIE TEMPS RÉEL & ALERTES (LIGHT) ====================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, c_light)
    add_header(slide7, "Messagerie Temps Réel & Alertes Sonores", "Stabilisation et robustesse de la connexion et de l'expérience utilisateur")
    
    left_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "STABILISATION TECHNIQUE"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark
    p.space_after = Pt(14)
    bullets = [
        "Configuration de Socket.io avec fallback de transport hybride : ['polling', 'websocket'].",
        "Contournement des pare-feu et blocages de tunnels en développement (Ngrok / Localtunnel).",
        "Injection automatique des en-têtes Ngrok pour bypasser la page d'avertissement de tunnel.",
        "Résolution des bugs d'ID : Cast strict des variables ID en String(id) pour le dispatch exact des messages.",
        "Refactoring du hook useSocket partagé et réutilisé."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    right_box = slide7.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "EXPÉRIENCE SONORE (EXPO-AV)"
    p_r.font.name = 'Segoe UI'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_orange
    p_r.space_after = Pt(14)
    solutions = [
        "Installation et configuration d'expo-av dans les applications mobiles.",
        "Création de l'utilitaire 'notificationSound.ts' pour zentraliser la lecture.",
        "Deux scénarios intégrés avec des fichiers audio MP3 légers :",
        "  - Ding-ding à la réception d'un nouveau message (sauf pour ses propres messages).",
        "  - Son arpège ascendant lors de la création d'annonces ou alertes.",
        "Nettoyage mémoire strict : déchargement systématique des fichiers audio après chaque lecture pour éviter tout gaspillage mémoire."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Segoe UI'
        p_s.font.size = Pt(13)
        p_s.font.color.rgb = c_gray_dark
        p_s.space_before = Pt(6)

    add_footer(slide7, 7, total_slides)

    # ==================== SLIDE 8 : AUTOMATISATION BACKEND & MODÉRATION (LIGHT) ====================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, c_light)
    add_header(slide8, "Automatisation Backend & Modération", "Système autonome pour réguler l'activité et assurer la conformité")
    
    left_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "EXPIRATION AUTOMATIQUE (72H)"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_sky
    p.space_after = Pt(14)
    bullets = [
        "Tâche de fond du serveur (`setInterval` toutes les 60s) exécutant la fonction d'expiration.",
        "Passage automatique des visites 'EN_ATTENTE' de plus de 72 heures au statut 'EXPIREE' / 'REFUSEE'.",
        "Enregistrement de la notification d'expiration persistante en BDD.",
        "Émission instantanée d'un signal WebSockets au client locataire concerné pour mettre à jour son UI en temps réel.",
        "Alerte visuelle (Alert.alert) et signal sonore (expo-av) immédiat déclenché chez le locataire."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    right_box = slide8.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "MODÉRATION SOUS 24H & EFFETS"
    p_r.font.name = 'Segoe UI'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_red
    p_r.space_after = Pt(14)
    solutions = [
        "Possibilité pour un locataire de signaler une annonce suspecte.",
        "L'annonce passe instantanément en statut 'SUSPENDUE' et disparaît de la recherche publique.",
        "Verrouillage total de la messagerie : le propriétaire ne peut plus envoyer de messages dans cette conversation.",
        "Bandeau rouge d'alerte de modération affiché sur le chat du propriétaire.",
        "Le propriétaire a 24h pour éditer et mettre à jour le bien pour soumission, dans la limite de 20 biens."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Segoe UI'
        p_s.font.size = Pt(13)
        p_s.font.color.rgb = c_gray_dark
        p_s.space_before = Pt(8)

    add_footer(slide8, 8, total_slides)

    # ==================== SLIDE 9 : SCÉNARIO DE DÉMO DE SOUTENANCE (LIGHT) ====================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, c_light)
    add_header(slide9, "Scénario de Démonstration Live", "Structure de la démonstration à réaliser devant les responsables")
    
    # Draw a custom horizontal process path with cards
    step_width = 3.6
    step_height = 2.1
    
    # Row 1 (Steps 1, 2, 3)
    # Step 1
    create_card(slide9, 0.8, 1.8, step_width, step_height, c_card_bg, c_sky)
    box = slide9.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(step_width - 0.2), Inches(step_height - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "1. Inscription & Publication"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_sky
    bullets = ["Connexion propriétaire", "Saisie d'un bien avec autocomplétion d'adresse Nominatim", "Le bien est publié en BDD"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "- " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Step 2
    create_card(slide9, 4.8, 1.8, step_width, step_height, c_card_bg, c_orange)
    box = slide9.shapes.add_textbox(Inches(4.9), Inches(1.9), Inches(step_width - 0.2), Inches(step_height - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2. Recherche & Itinéraire"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_orange
    bullets = ["Connexion locataire", "Sélection du bien sur la carte", "Tracé d'itinéraire OSRM (fallback Haversine si hors-ligne)"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "- " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Step 3
    create_card(slide9, 8.8, 1.8, step_width, step_height, c_card_bg, c_dark)
    box = slide9.shapes.add_textbox(Inches(8.9), Inches(1.9), Inches(step_width - 0.2), Inches(step_height - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "3. Demande de Visite"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_dark
    bullets = ["Clic sur bouton 'Contacter' bloqué", "Choix de la date et de l'heure", "Envoi d'une demande 'EN_ATTENTE'"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "- " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Row 2 (Steps 4, 5, 6)
    # Step 4
    create_card(slide9, 0.8, 4.3, step_width, step_height, c_card_bg, c_red)
    box = slide9.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(step_width - 0.2), Inches(step_height - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "4. Validation & Chat"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_red
    bullets = ["Propriétaire accepte la visite", "Le salon de discussion s'ouvre", "Échange en temps réel (sons expo-av)"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "- " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Step 5
    create_card(slide9, 4.8, 4.3, step_width, step_height, c_card_bg, c_gray_dark)
    box = slide9.shapes.add_textbox(Inches(4.9), Inches(4.4), Inches(step_width - 0.2), Inches(step_height - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "5. Expiration 72h (Simulée)"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_gray_dark
    bullets = ["Visite non validée expire en BDD", "Signal socket instantané", "Alerte visuelle et sonore chez le locataire"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "- " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    # Step 6
    create_card(slide9, 8.8, 4.3, step_width, step_height, c_card_bg, c_orange)
    box = slide9.shapes.add_textbox(Inches(8.9), Inches(4.4), Inches(step_width - 0.2), Inches(step_height - 0.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "6. Modération & Suspension"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = c_orange
    bullets = ["Signalement par le locataire", "L'annonce passe à SUSPENDUE", "Verrouillage de la messagerie du propriétaire"]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = "- " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(2)

    add_footer(slide9, 9, total_slides)

    # ==================== SLIDE 10 : PERFORMANCE, SÉCURITÉ & PERSPECTIVES (LIGHT) ====================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, c_light)
    add_header(slide10, "Sécurité, Robustesse & Perspectives", "Garantir la fiabilité en production et préparer le passage à l'échelle")
    
    left_box = slide10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "SÉCURITÉ ET ROBUSTESSE DU CODE"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark
    p.space_after = Pt(14)
    bullets = [
        "Hachage sécurisé Bcrypt en base de données et authentification via tokens JWT.",
        "Typage rigoureux en TypeScript vérifié avec 'npx tsc --noEmit' sur l'ensemble du projet.",
        "Gestion efficace des ressources mobiles : libération immédiate de la mémoire de lecture d'expo-av après chaque bip.",
        "Résolution des problèmes d'importation de modules natifs sur le web (wrappers conditionnels pour maps)."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = c_gray_dark
        p_b.space_before = Pt(8)

    right_box = slide10.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "PERSPECTIVES DE PRODUCTION"
    p_r.font.name = 'Segoe UI'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_sky
    p_r.space_after = Pt(14)
    solutions = [
        "Migration des serveurs de développement vers un hébergement stable (Frontend Vercel, Backend Render, PostgreSQL Neon DB).",
        "Remplacement des alertes Socket.io en tâche de fond par de vraies notifications push via Expo Notifications, Google FCM et Apple APNs (nécessite l'attribution d'un projectId Expo Cloud).",
        "Déploiement des applications mobiles sur l'App Store et Google Play Store.",
        "Intégration d'un module de paiement mobile Money (Wave, Orange, MTN) pour les réservations."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Segoe UI'
        p_s.font.size = Pt(13)
        p_s.font.color.rgb = c_gray_dark
        p_s.space_before = Pt(8)

    add_footer(slide10, 10, total_slides)

    # ==================== SLIDE 11 : CONCLUSION (DARK) ====================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, c_dark)
    
    # Accent bar on left
    accent_bar = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_orange
    accent_bar.line.color.rgb = c_orange

    title_box = slide11.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_brand = tf.paragraphs[0]
    p_brand.text = "ATTOUHOME  •  CONCLUSION"
    p_brand.font.name = 'Segoe UI'
    p_brand.font.size = Pt(16)
    p_brand.font.bold = True
    p_brand.font.color.rgb = c_orange
    p_brand.space_after = Pt(12)
    
    p_title = tf.add_paragraph()
    p_title.text = "Une Plateforme Prête pour sa Démonstration"
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.space_after = Pt(18)
    
    points = [
        "✔ Solution complète et intégrée (2 applications mobiles et 1 tableau de bord d'administration).",
        "✔ Fonctionnalités clés stabilisées : Itinéraires OSRM, Chat Socket.io et alertes audio.",
        "✔ Robustesse opérationnelle grâce à l'automatisation backend (expiration et modération).",
        "✔ Architecture typée TypeScript, extensible et prête pour un déploiement cloud d'envergure."
    ]
    for pt in points:
        p_pt = tf.add_paragraph()
        p_pt.text = pt
        p_pt.font.name = 'Segoe UI'
        p_pt.font.size = Pt(16)
        p_pt.font.bold = True
        p_pt.font.color.rgb = c_sky
        p_pt.space_before = Pt(8)
        
    p_thanks = tf.add_paragraph()
    p_thanks.text = "Merci pour votre attention. Place à la démonstration en direct !"
    p_thanks.font.name = 'Segoe UI'
    p_thanks.font.size = Pt(18)
    p_thanks.font.italic = True
    p_thanks.font.bold = True
    p_thanks.font.color.rgb = c_white
    p_thanks.space_before = Pt(30)

    add_footer(slide11, 11, total_slides, is_dark=True)
    
    # Enregistrer la présentation
    file_path = "Soutenance_Projet_AttouHome.pptx"
    prs.save(file_path)
    print(f"Présentation PowerPoint générée avec succès : {file_path}")

if __name__ == '__main__':
    create_soutenance_presentation()
