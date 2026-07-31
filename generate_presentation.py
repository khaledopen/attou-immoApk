import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Configuration des dimensions (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Couleurs du thème AttouHome
    c_dark = RGBColor(15, 23, 42)      # #0f172a (Bleu nuit profond)
    c_sky = RGBColor(14, 165, 233)     # #0ea5e9 (Bleu ciel)
    c_orange = RGBColor(249, 115, 22)  # #f97316 (Orange vibrant)
    c_red = RGBColor(239, 68, 68)      # #ef4444 (Rouge Alerte)
    c_gray = RGBColor(71, 85, 105)     # #475569 (Gris ardoise)
    c_light = RGBColor(248, 250, 252)  # #f8fafc (Gris très clair)
    c_white = RGBColor(255, 255, 255)  # Blanc
    
    # Layouts standards (layout 6 est vide)
    blank_layout = prs.slide_layouts[6]
    
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, subtitle_text=None, is_dark=False):
        # Boîte de titre principal
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.margin_left = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = c_white if is_dark else c_dark
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Arial'
            p2.font.size = Pt(14)
            p2.font.italic = True
            p2.font.color.rgb = c_sky
            p2.space_before = Pt(6)

    # ==================== SLIDE 1 : TITRE (DARK THEME) ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, c_dark)
    
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_brand = tf.paragraphs[0]
    p_brand.text = "ATTOUHOME  •  BILAN TECHNIQUE DES ÉVOLUTIONS"
    p_brand.font.name = 'Arial'
    p_brand.font.size = Pt(14)
    p_brand.font.bold = True
    p_brand.font.color.rgb = c_orange
    p_brand.space_after = Pt(20)
    
    p_title = tf.add_paragraph()
    p_title.text = "Évolutions de la Plateforme\nMessagerie, Alertes, Visites & Modération"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(44)
    p_title.font.bold = True
    p_title.font.color.rgb = c_white
    p_title.space_after = Pt(18)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "Présentation complète des fonctionnalités développées et stabilisées"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = c_sky
    p_sub.space_before = Pt(10)

    # ==================== SLIDE 2 : MESSAGERIE TEMPS RÉEL ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, c_light)
    add_header(slide2, "1. Messagerie Temps Réel Résiliente", "Stabilisation et robustesse de la connexion client-serveur")
    
    left_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "LE PROBLÈME RENCONTRÉ"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(14)
    bullets = [
        "Échec de réception des messages en temps réel sur les appareils mobiles.",
        "Le tunnel Ngrok bloquait les connexions WebSockets pures.",
        "Routage restreint par transports forcés ('websocket' uniquement).",
        "Échecs de comparaison d'ID dus à des différences de type (String vs Object)."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = c_gray
        p_b.space_before = Pt(8)
        
    right_box = slide2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "RÉSOLUTIONS APPORTÉES"
    p_r.font.name = 'Arial'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_sky
    p_r.space_after = Pt(14)
    solutions = [
        "Ajout du fallback de transport : ['polling', 'websocket'] pour contourner les blocages réseau.",
        "Injection automatique des en-têtes Ngrok pour bypasser la page d'avertissement de tunnel.",
        "Standardisation globale du hook useSocket et des connexions dans l'owner-app et la tenant-app.",
        "Cast explicite String(id) pour sécuriser la logique de dispatch des messages reçus."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(14)
        p_s.font.color.rgb = c_gray
        p_s.space_before = Pt(8)

    # ==================== SLIDE 3 : SONS DE NOTIFICATION ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, c_light)
    add_header(slide3, "2. Alertes Sonores (expo-av)", "Création d'une identité sonore pour une réactivité maximale")
    
    left_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "INTÉGRATION TECHNIQUE"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark
    p.space_after = Pt(14)
    bullets = [
        "Installation de la librairie native Expo 'expo-av' dans les deux applications mobiles.",
        "Création d'un utilitaire partagé 'notificationSound.ts' gérant la lecture et le déchargement.",
        "Fichiers audios MP3 légers créés et intégrés directement sous 'assets/sounds/'."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = c_gray
        p_b.space_before = Pt(8)
        
    right_box = slide3.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "SCÉNARIOS AUDIOS INTÉGRÉS"
    p_r.font.name = 'Arial'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_orange
    p_r.space_after = Pt(14)
    solutions = [
        "Réception d'un message : Un son court doux type 'ding-ding' (ne se déclenche pas pour ses propres messages).",
        "Création d'un bien / Notifications : Un son arpège ascendant distinctif type 'notification'.",
        "Nettoyage automatique : Déchargement de la mémoire après chaque lecture pour éviter les fuites audio.",
        "Intégration globale dans les layouts mobiles pour capter toutes les alertes en tâche de fond."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(14)
        p_s.font.color.rgb = c_gray
        p_s.space_before = Pt(8)

    # ==================== SLIDE 4 : DEMANDE DE VISITE DATE ET HEURE ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, c_light)
    add_header(slide4, "3. Flux de Demande de Visite (Date & Heure)", "Double validation intuitive et prise de rendez-vous complète")
    
    left_box = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "EXPÉRIENCE LOCATAIRE (TENANT)"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark
    p.space_after = Pt(14)
    bullets = [
        "La carte propriétaire et l'avatar sur la fiche du bien sont désormais cliquables pour contacter le propriétaire.",
        "Alerte si aucune visite n'est encore planifiée pour rediriger vers la demande préalable.",
        "Enchaînement automatique : Sélection du jour sur le calendrier, puis ouverture immédiate du sélecteur d'heure (TimePicker).",
        "Combinaison robuste de la date et de l'heure locale converties en ISO avant envoi au serveur."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = c_gray
        p_b.space_before = Pt(8)
        
    right_box = slide4.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "EXPÉRIENCE PROPRIÉTAIRE (OWNER)"
    p_r.font.name = 'Arial'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_sky
    p_r.space_after = Pt(14)
    solutions = [
        "Affichage clair de la date et de l'heure sur l'écran d'accueil/Dashboard (ex: '4 juin à 15:30').",
        "Rendu complet dans la liste des visites pour faciliter la planification de l'agenda.",
        "Vérification instantanée lors du chargement des demandes."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(14)
        p_s.font.color.rgb = c_gray
        p_s.space_before = Pt(8)

    # ==================== SLIDE 5 : EXPIRATION AUTOMATIQUE ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, c_light)
    add_header(slide5, "4. Expiration Automatique sous 72 Heures", "Gestion dynamique et notification automatique des rendez-vous échus")
    
    left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "RÈGLE ET LOGIQUE BACKEND"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_dark
    p.space_after = Pt(14)
    bullets = [
        "Une tâche de fond du serveur (`setInterval` de 60 secondes) exécute la routine `expireVisits`.",
        "Toutes les demandes en statut 'EN_ATTENTE' datant de plus de 72h passent automatiquement en statut 'REFUSEE'.",
        "Création automatique d'une notification persistante dans la base de données PostgreSQL."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = c_gray
        p_b.space_before = Pt(8)
        
    right_box = slide5.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "NOTIFICATIONS TEMPS RÉEL CLIENT"
    p_r.font.name = 'Arial'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_orange
    p_r.space_after = Pt(14)
    solutions = [
        "Émission socket d'une notification instantanée 'VISITE_EXPIREE' au locataire concerné.",
        "Le client mobile intercepte l'événement, déclenche une alerte système (Alert.alert) et joue immédiatement le son de notification pour avertir l'utilisateur.",
        "Mise à jour directe du statut dans l'historique des visites du locataire."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(14)
        p_s.font.color.rgb = c_gray
        p_s.space_before = Pt(8)

    # ==================== SLIDE 6 : MODÉRATION ET REMISE EN LOCATION ====================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, c_light)
    add_header(slide6, "5. Signalement & Modération sous 24h", "Sécurité des annonces et contrôle de la messagerie")
    
    left_box = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    p = ltf.paragraphs[0]
    p.text = "SIGNALEMENT ET EFFET DIRECT"
    p.font.name = 'Arial'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(14)
    bullets = [
        "Le signalement passe l'annonce en 'SUSPENDUE' et la retire des recherches.",
        "Délai officiel de 24h donné au propriétaire pour corriger les erreurs (photo incorrecte, fausse description, etc.).",
        "Remise en location sécurisée : validation renforcée de la limite (20 annonces max)."
    ]
    for b in bullets:
        p_b = ltf.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = c_gray
        p_b.space_before = Pt(8)
        
    right_box = slide6.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    p_r = rtf.paragraphs[0]
    p_r.text = "CONTRÔLE DE MESSAGERIE ET BANNIÈRE"
    p_r.font.name = 'Arial'
    p_r.font.size = Pt(16)
    p_r.font.bold = True
    p_r.font.color.rgb = c_dark
    p_r.space_after = Pt(14)
    solutions = [
        "Affichage d'un badge rouge 'Modération' et d'une bannière d'alerte officielle de l'administration dans le chat du propriétaire.",
        "Verrouillage total de la messagerie : le propriétaire ne peut plus envoyer de messages dans cette conversation.",
        "Mise à jour instantanée des listes de biens chez le locataire via focus de page et sockets."
    ]
    for s in solutions:
        p_s = rtf.add_paragraph()
        p_s.text = "✔ " + s
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(14)
        p_s.font.color.rgb = c_gray
        p_s.space_before = Pt(8)

    # ==================== SLIDE 7 : CONCLUSION (DARK THEME) ====================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, c_dark)
    
    title_box = slide7.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p_conc = tf.paragraphs[0]
    p_conc.text = "BILAN DE LA VALEUR AJOUTÉE"
    p_conc.font.name = 'Arial'
    p_conc.font.size = Pt(16)
    p_conc.font.bold = True
    p_conc.font.color.rgb = c_orange
    p_conc.space_after = Pt(20)
    
    points = [
        "1. Expérience client de haute qualité : Des retours sonores immersifs et un sélecteur de date/heure très fluide.",
        "2. Sécurisation commerciale : Pas de chat possible entre locataires et propriétaires sans demande de visite validée au préalable.",
        "3. Automatisation intelligente : Système d'expiration des visites à 72h et alertes de modération à 24h gérés automatiquement par le serveur.",
        "4. Résilience technique : Connexions WebSocket stabilisées avec fallback HTTP polling en cas d'utilisation de tunnels Ngrok."
    ]
    for pt in points:
        p_pt = tf.add_paragraph()
        p_pt.text = pt
        p_pt.font.name = 'Arial'
        p_pt.font.size = Pt(20)
        p_pt.font.bold = True
        p_pt.font.color.rgb = c_white
        p_pt.space_before = Pt(12)
        
    prs.save("Bilan_Optimisations_AttouHome.pptx")
    print("Présentation PowerPoint générée avec succès : Bilan_Optimisations_AttouHome.pptx")

if __name__ == '__main__':
    create_presentation()
